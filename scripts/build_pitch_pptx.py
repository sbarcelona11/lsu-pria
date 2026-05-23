from __future__ import annotations

import argparse
import json
import re
from pathlib import Path

from _bootstrap import ensure_repo_root_on_path

ensure_repo_root_on_path()

from pptx import Presentation  # type: ignore
from pptx.dml.color import RGBColor  # type: ignore
from pptx.enum.text import PP_ALIGN  # type: ignore
from pptx.util import Inches, Pt  # type: ignore


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser()
    p.add_argument("--config", default="deliverables/config.json")
    p.add_argument("--pitch-md", default="deliverables/out/_tmp/pitch_resolved.md")
    p.add_argument("--out", default="deliverables/out/pitch.pptx")
    return p.parse_args()


def _read(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def _apply_placeholders(md: str, cfg: dict) -> str:
    members = cfg.get("members", [])
    if isinstance(members, list):
        members_s = ", ".join(str(x) for x in members)
    else:
        members_s = str(members)
    mapping = {
        "{{TITLE}}": str(cfg.get("title", "VC-pria")),
        "{{COURSE}}": str(cfg.get("course", "")),
        "{{GROUP}}": str(cfg.get("group", "")),
        "{{MEMBERS}}": members_s,
        "{{DATE}}": str(cfg.get("date", "")),
        "{{PITCH_MAX_MIN}}": str(cfg.get("target_pitch_minutes", 8)),
    }
    out = md
    for k, v in mapping.items():
        out = out.replace(k, v)
    return out


def _md_to_lines(md: str) -> list[tuple[str, str]]:
    out: list[tuple[str, str]] = []
    for raw in md.splitlines():
        line = raw.rstrip()
        if not line.strip():
            out.append(("blank", ""))
            continue
        if line.startswith("### "):
            out.append(("h3", line[4:].strip()))
        elif line.startswith("## "):
            out.append(("h2", line[3:].strip()))
        elif line.startswith("# "):
            out.append(("h1", line[2:].strip()))
        elif line.lstrip().startswith("- "):
            out.append(("bullet", line.lstrip()[2:].strip()))
        else:
            out.append(("para", line.strip()))
    return out


def _strip_md(s: str) -> str:
    s = re.sub(r"[*_`]", "", s)
    s = s.replace("—", "-").replace("–", "-").replace("\u00a0", " ")
    return s.strip()


def _parse_pitch_slides(md: str) -> list[tuple[str, list[str]]]:
    slides: list[tuple[str, list[str]]] = []
    current_title = "Pitch"
    current_body: list[str] = []
    for kind, text in _md_to_lines(md):
        if kind == "h2" and text.lower().startswith("slide"):
            if current_body or current_title:
                slides.append((current_title, current_body))
            current_title = text
            current_body = []
        elif kind in ("bullet", "para"):
            current_body.append(text)
    if current_body:
        slides.append((current_title, current_body))
    slides = [s for s in slides if s[0] and (s[1] or "slide" in s[0].lower())]
    return slides


def _title_display(title: str) -> str:
    t = title.strip()
    if "—" in t:
        t = t.split("—", 1)[1].strip()
    elif "-" in t and t.lower().startswith("slide"):
        parts = t.split("-", 1)
        if len(parts) == 2:
            t = parts[1].strip()
    return t or "Slide"


def main() -> None:
    args = parse_args()
    repo = Path(__file__).resolve().parents[1]
    cfg_path = repo / args.config
    cfg = json.loads(cfg_path.read_text(encoding="utf-8")) if cfg_path.exists() else {}

    pitch_md_path = repo / args.pitch_md
    if not pitch_md_path.exists():
        # Fallback to source pitch.md if resolved one doesn't exist yet.
        pitch_md_path = repo / "deliverables" / "pitch.md"
        md = _apply_placeholders(_read(pitch_md_path), cfg)
    else:
        md = _read(pitch_md_path)

    slides = _parse_pitch_slides(md)
    if not slides:
        slides = [("Slide 1 - Pitch", ["(sin contenido)"])]

    prs = Presentation()
    # Force 16:9 widescreen.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG = RGBColor(0x0B, 0x10, 0x20)  # #0b1020
    BAR = RGBColor(0x1E, 0x28, 0x50)  # #1e2850
    FG = RGBColor(0xE8, 0xEC, 0xFF)  # #e8ecff
    FG_DIM = RGBColor(0xA0, 0xA6, 0xC6)  # #a0a6c6

    blank_layout = prs.slide_layouts[6]
    total = len(slides)
    for idx, (title, body) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank_layout)

        # Background fill
        bg = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE, avoid importing enum for compat
            0,
            0,
            prs.slide_width,
            prs.slide_height,
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG
        bg.line.fill.background()

        # Header bar
        header = slide.shapes.add_shape(
            1,
            Inches(0.45),
            Inches(0.4),
            prs.slide_width - Inches(0.9),
            Inches(0.85),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = BAR
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.7),
            Inches(0.5),
            prs.slide_width - Inches(1.4),
            Inches(0.65),
        )
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = _strip_md(_title_display(title))
        run.font.size = Pt(34)
        run.font.bold = True
        run.font.color.rgb = FG
        p.alignment = PP_ALIGN.LEFT

        # Bullets
        bullets = [_strip_md(x) for x in body if _strip_md(x)]
        if len(bullets) > 8:
            bullets = bullets[:8] + ["(ver detalle en informe)"]

        body_box = slide.shapes.add_textbox(
            Inches(0.95),
            Inches(1.55),
            prs.slide_width - Inches(1.6),
            prs.slide_height - Inches(2.2),
        )
        btf = body_box.text_frame
        btf.word_wrap = True
        btf.clear()
        for i, b in enumerate(bullets):
            if i == 0:
                p = btf.paragraphs[0]
            else:
                p = btf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(26)
            p.font.color.rgb = FG

        # Footer (left)
        footer_l = slide.shapes.add_textbox(
            Inches(0.45),
            prs.slide_height - Inches(0.45),
            Inches(8.0),
            Inches(0.3),
        )
        ftf = footer_l.text_frame
        ftf.clear()
        p = ftf.paragraphs[0]
        p.text = "VC-pria — Pitch"
        p.font.size = Pt(14)
        p.font.color.rgb = FG_DIM

        # Footer (right)
        footer_r = slide.shapes.add_textbox(
            prs.slide_width - Inches(2.0),
            prs.slide_height - Inches(0.45),
            Inches(1.55),
            Inches(0.3),
        )
        rtf = footer_r.text_frame
        rtf.clear()
        p = rtf.paragraphs[0]
        p.text = f"{idx}/{total}"
        p.font.size = Pt(14)
        p.font.color.rgb = FG_DIM
        p.alignment = PP_ALIGN.RIGHT

    out_path = (repo / args.out).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Wrote: {out_path}")


if __name__ == "__main__":
    main()

